# This script uses a class-based approach to create an Excel report,
# add data and a native chart to it, and then generate an image of that
# chart for potential use in other applications like PowerPoint.
#
# It requires the following libraries to be installed:
# pip install openpyxl
# pip install python-pptx
# pip install matplotlib

import os
import openpyxl,pptx
from openpyxl.workbook import Workbook
from openpyxl.chart import ScatterChart, Reference
from openpyxl.chart.series_factory import SeriesFactory as Series
from openpyxl.chart.shapes import GraphicalProperties
from openpyxl.drawing.line import LineProperties
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

class ExcelReport:
    def __init__(self):
        self.path = os.getcwd()
        self.Report_Path = os.path.join(self.path, "Report")
        if not os.path.exists(self.Report_Path):
            os.makedirs(self.Report_Path)
        os.chdir(self.Report_Path)

        self.workbook = Workbook()
        self.results = self.workbook.active
        self.ss_sheet1 = self.workbook['Sheet']
        self.ss_sheet1.title = 'Test results'
        self.row = 1

    def add_headers(self, headers):
        for col_idx, header_title in enumerate(headers, 1):
            self.results.cell(row=self.row, column=col_idx, value=header_title)

    def add_values(self, values, filename):
        self.row += 1
        for col_idx, value in enumerate(values, 1):
            self.results.cell(row=self.row, column=col_idx, value=value)
        self.workbook.save(filename)

    def add_chart_to_excel(self, filename):

        print(f"\nAdding chart to '{filename}'.")
        try:
            workbook = openpyxl.load_workbook(filename)
            sheet = workbook.active

            chart = ScatterChart()
            chart.title = "Parameter Analysis"
            chart.style = 13
            chart.height = 10
            chart.width = 20
            chart.y_axis.title = "Parameter Values"
            chart.x_axis.title = "Timestamp (s)"
            chart.legend.position = 't'

            # Define data ranges based on the class structure (headers in row 1, data from row 2)
            max_data_row = sheet.max_row

            # --- Series 1 ---
            x_values = Reference(sheet, min_col=1, min_row=2, max_row=max_data_row)
            y_values1 = Reference(sheet, min_col=2, min_row=2, max_row=max_data_row)
            series1 = Series(values=y_values1, xvalues=x_values, title_from_data=False,
                             title=sheet.cell(row=1, column=2).value)
            series1.spPr = GraphicalProperties(ln=LineProperties())
            chart.series.append(series1)

            # --- Series 2 ---
            y_values2 = Reference(sheet, min_col=3, min_row=2, max_row=max_data_row)
            series2 = Series(values=y_values2, xvalues=x_values, title_from_data=False,
                             title=sheet.cell(row=1, column=3).value)
            series2.spPr = GraphicalProperties(ln=LineProperties())
            chart.series.append(series2)

            # Add the chart to the worksheet
            sheet.add_chart(chart, "F2")

            workbook.save(filename)
            print("Chart added to Excel file successfully.")

        except FileNotFoundError:
            print(f"Error: The file '{filename}' was not found.")

    def validate_data_in_excel(self,filename, column_index=3, limit=20):
        """
        Reads an Excel file, checks if values in a specific column are within a limit,
        and returns messages for any deviations.
        """
        print("\n--- Starting Data Validation ---")
        out_of_limit_messages = []
        try:
            workbook = openpyxl.load_workbook(filename)
            sheet = workbook.active

            for row in range(2, sheet.max_row + 1):
                cell_value = sheet.cell(row=row, column=column_index).value
                Parameter1 = sheet.cell(row=row, column=1).value
                Parameter3 = sheet.cell(row=row, column=3).value

                if cell_value is not None and cell_value > limit:
                    message = f"- Parameter 1 {Parameter1} with the parameter 3 {Parameter3} is out of limit with {cell_value} ."
                    print(f"ALERT: Deviation found for {Parameter1}.")
                    out_of_limit_messages.append(message)

            if not out_of_limit_messages:
                message ="All data is within the specified limit."
                print("All data is within the specified limit.")
                out_of_limit_messages.append(message)

            print("--- Data Validation Complete ---")
            return out_of_limit_messages

        except FileNotFoundError:
            print(f"Error: The file '{filename}' was not found.")
            return []

    def create_chart_as_image(self, excel_filename, output_image_filename="chart.png"):
        """
        Reads data from Excel and creates a multi-series plot as a saved image file.
        This method is designed to work with the data structure of the ExcelReport class.
        """
        print(f"\nGenerating chart image from '{excel_filename}'...")
        try:
            workbook = openpyxl.load_workbook(excel_filename)
            sheet = workbook.active

            # --- Dynamically Read Headers for Labels (from row 1) ---
            header_row = 1
            x_label = sheet.cell(row=header_row, column=1).value
            y1_label = sheet.cell(row=header_row, column=2).value
            y2_label = sheet.cell(row=header_row, column=3).value

            # --- Read Data (from row 2 onwards) ---
            x_data, y1_data, y2_data = [], [], []
            for row in range(2, sheet.max_row + 1):
                x_data.append(sheet.cell(row=row, column=1).value)
                y1_data.append(sheet.cell(row=row, column=2).value)
                y2_data.append(sheet.cell(row=row, column=3).value)

            plt.figure(figsize=(10, 6))

            # Plot both data series using the dynamically read labels
            plt.plot(x_data, y1_data, marker='o', linestyle='-', color='coral', markersize=8, label=y1_label)
            plt.plot(x_data, y2_data, marker='o', linestyle='-', color='skyblue', markersize=8, label=y2_label)

            plt.title('Parameter Analysis', fontsize=16)
            plt.xlabel(x_label, fontsize=12)
            plt.ylabel('Values', fontsize=12)
            plt.grid(True, which='both', linestyle='--', linewidth=0.5)
            plt.legend()
            plt.tight_layout()

            plt.savefig(output_image_filename)
            plt.close()
            print(f"Chart saved as '{output_image_filename}'.")
            return output_image_filename

        except FileNotFoundError:
            print(f"Error: The file '{excel_filename}' was not found.")
            return None

    def export_data_to_powerpoint(self,excel_filename, ppt_filename, chart_image_path=None, deviation_messages=None):
        """
        Creates a PowerPoint with a data table, and conditionally adds a slide with a chart and notes.
        """
        print(f"\nExporting data from '{excel_filename}' to '{ppt_filename}'...")
        try:
            workbook = openpyxl.load_workbook(excel_filename)
            sheet = workbook.active
            prs = Presentation()

            # --- Slide 1: Data Table ---
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Report"

            data_for_ppt = [tuple(cell.value for cell in row) for row in sheet.iter_rows()]
            rows, cols = len(data_for_ppt), len(data_for_ppt[0])
            table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2.0), Inches(8), Inches(0.8) * rows)
            table = table_shape.table

            for r_idx, row_data in enumerate(data_for_ppt):
                for c_idx, cell_data in enumerate(row_data):
                    table.cell(r_idx, c_idx).text = str(cell_data)

            # --- Slide 2: Conditional Chart and Deviations ---
            if chart_image_path and deviation_messages:
                print("Adding deviation analysis slide to PowerPoint...")
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Analysis"

                pic = slide.shapes.add_picture(chart_image_path, Inches(0.5), Inches(1.5), width=Inches(6))

                textbox_left = Inches(1)
                textbox_top = pic.top + pic.height + Inches(0.25)  # Position below the picture with a small gap
                textbox_width = Inches(8)
                textbox_height = Inches(1.5)

                textbox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)

                # textbox = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(3), Inches(4))
                text_frame = textbox.text_frame
                text_frame.text = "Data Limit Breaches:"
                p = text_frame.add_paragraph()
                p.text = "\n".join(deviation_messages)
                p.font.size = pptx.util.Pt(14)
            else:
                print("NOT ABLE TO ADD")

            prs.save(ppt_filename)
            print(f"Presentation '{ppt_filename}' created successfully.")

        except FileNotFoundError:
            print(f"Error: The file '{excel_filename}' was not found.")

if __name__ == "__main__":
    Excel = ExcelReport()
    filename = "Report_param.xlsx"

    # Define headers and add them to the first row
    header = ["Parameter1", "Parameter2", "Parameter3","Parameter4", ]
    Excel.add_headers(header)
    values = [5.4,6.3,7.1,8.3]

    # Add 20 rows of sample data
    for i in range(5):
        # Make data more interesting for a scatter plot
        values[0] = i + 1
        values[1] = i + 1
        values[2] = i + 5
        values[3] = i + 1
        Excel.add_values(list(values),filename)

    Excel.add_chart_to_excel(filename)

    deviation_messages = Excel.validate_data_in_excel(filename, column_index=3, limit=9)
    chart_image = None
    if deviation_messages:
        chart_image = Excel.create_chart_as_image(filename)

    Excel.export_data_to_powerpoint(filename,"Report_param.pptx",chart_image,deviation_messages)

